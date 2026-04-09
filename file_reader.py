# -*- coding: utf-8 -*-
"""
文件读取器模块 - 支持多种文件格式
"""
import os
import mmap
import chardet
from constants import CHUNK_SIZE, LARGE_FILE_THRESHOLD

# 尝试导入各种文档处理库
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    import xlrd
    XLS_AVAILABLE = True
except ImportError:
    XLS_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# 尝试导入流量包处理库
try:
    from scapy.all import rdpcap, Raw, IP, TCP, UDP, DNS
    # 尝试导入 HTTP 相关类
    try:
        from scapy.layers.http import HTTPRequest, HTTPResponse
        HTTP_AVAILABLE = True
    except ImportError:
        HTTP_AVAILABLE = False
        HTTPRequest = None
        HTTPResponse = None
    SCAPY_AVAILABLE = True
except ImportError:
    SCAPY_AVAILABLE = False
    HTTP_AVAILABLE = False
    HTTPRequest = None
    HTTPResponse = None

# 尝试导入 Windows 事件日志处理库
try:
    from Evtx.Evtx import Evtx
    import xml.etree.ElementTree as ET
    EVTX_AVAILABLE = True
except ImportError:
    try:
        # 备用导入方式
        import Evtx.Evtx as evtx_module
        Evtx = evtx_module.Evtx
        EVTX_AVAILABLE = True
    except ImportError:
        EVTX_AVAILABLE = False
        Evtx = None


class FileReader:
    """支持多种文件格式的读取器"""
    
    def __init__(self, auto_detect_encoding=True):
        self.auto_detect_encoding = auto_detect_encoding
        self.encoding_cache = {}
    
    def get_file_extension(self, file_path):
        """获取文件扩展名（小写）"""
        return os.path.splitext(file_path)[1].lower()
    
    def is_text_file(self, file_path):
        """判断是否为纯文本文件"""
        text_extensions = [
            '.txt', '.log', '.csv', '.json', '.xml', '.html', '.htm',
            '.js', '.jsx', '.ts', '.tsx', '.mjs', '.cjs',
            '.py', '.pyw', '.java', '.c', '.cpp', '.cc', '.cxx', '.h', '.hpp',
            '.cs', '.php', '.rb', '.go', '.rs', '.swift', '.kt', '.scala',
            '.sql', '.sh', '.bash', '.bat', '.cmd', '.ps1',
            '.css', '.scss', '.sass', '.less',
            '.md', '.markdown', '.rst',
            '.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf',
            '.properties', '.gradle', '.maven',
            '.vue', '.jsx', '.tsx',
            '.dockerfile', '.makefile',
            '.gitignore', '.gitattributes',
            '.env', '.env.local'
        ]
        return self.get_file_extension(file_path) in text_extensions
    
    def detect_encoding(self, file_path):
        """检测文件编码"""
        if file_path in self.encoding_cache:
            return self.encoding_cache[file_path]

        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10240)
                if not raw_data.strip():
                    return 'utf-8'
                result = chardet.detect(raw_data)
                encoding = result['encoding'] or 'utf-8'
                encoding = encoding.lower().replace('utf-16le', 'utf-16').replace('utf-16be', 'utf-16')
                try:
                    raw_data.decode(encoding)
                    self.encoding_cache[file_path] = encoding
                    return encoding
                except UnicodeDecodeError:
                    pass

            for enc in ['utf-8', 'gbk', 'gb18030', 'big5', 'utf-16']:
                try:
                    with open(file_path, 'r', encoding=enc) as f:
                        f.read(100)
                    self.encoding_cache[file_path] = enc
                    return enc
                except (UnicodeDecodeError, LookupError):
                    continue
            return 'utf-8'
        except Exception:
            return 'utf-8'
    
    def read_text_file(self, file_path):
        """读取纯文本文件（优化版本，使用mmap加速大文件读取）"""
        try:
            file_size = os.path.getsize(file_path)
            
            # 检查是否为二进制文件
            with open(file_path, 'rb') as f:
                head = f.read(min(1024, file_size))
                if b'\x00' in head:
                    return None

            if self.auto_detect_encoding:
                encoding = self.detect_encoding(file_path)
            else:
                encoding = 'utf-8'

            # 对于大文件，使用mmap内存映射加速读取
            if file_size > LARGE_FILE_THRESHOLD:
                try:
                    # 使用mmap加速大文件读取
                    with open(file_path, 'rb') as f:
                        with mmap.mmap(f.fileno(), 0, access=mmap.ACCESS_READ) as mm:
                            # 分块解码，避免一次性解码导致内存问题
                            chunks = []
                            chunk_size = CHUNK_SIZE * 8  # 进一步增大块大小（16MB），减少系统调用和内存分配
                            offset = 0
                            
                            while offset < len(mm):
                                chunk = mm[offset:offset + chunk_size]
                                if not chunk:
                                    break
                                
                                try:
                                    decoded = chunk.decode(encoding, errors='ignore')
                                    chunks.append(decoded)
                                except Exception:
                                    try:
                                        decoded = chunk.decode('utf-8', errors='ignore')
                                        chunks.append(decoded)
                                    except Exception:
                                        decoded = chunk.decode('gbk', errors='ignore')
                                        chunks.append(decoded)
                                
                                offset += chunk_size
                            
                            return ''.join(chunks)
                except (OSError, ValueError, MemoryError):
                    # mmap失败，回退到普通读取（使用更大的缓冲区）
                    chunks = []
                    with open(file_path, 'rb', buffering=CHUNK_SIZE * 16) as f:  # 32MB缓冲区
                        while True:
                            chunk = f.read(CHUNK_SIZE * 8)  # 16MB读取块
                            if not chunk:
                                break
                            try:
                                chunks.append(chunk.decode(encoding, errors='ignore'))
                            except Exception:
                                try:
                                    chunks.append(chunk.decode('utf-8', errors='ignore'))
                                except Exception:
                                    chunks.append(chunk.decode('gbk', errors='ignore'))
                    return ''.join(chunks)
            else:
                # 小文件，直接读取（使用较大缓冲区，提升IO性能）
                with open(file_path, 'rb', buffering=CHUNK_SIZE * 2) as f:  # 4MB缓冲区
                    data = f.read()
                try:
                    return data.decode(encoding, errors='ignore')
                except Exception:
                    try:
                        return data.decode('utf-8', errors='ignore')
                    except Exception:
                        return data.decode('gbk', errors='ignore')
        except Exception as e:
            print(f"读取文本文件 {file_path} 时出错: {str(e)}")
            return None
    
    def read_docx(self, file_path):
        """读取 DOCX 文件"""
        if not DOCX_AVAILABLE:
            return None
        
        try:
            doc = Document(file_path)
            text_parts = []
            
            # 读取段落
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # 读取表格
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            return "\n".join(text_parts)
        except Exception as e:
            print(f"读取 DOCX 文件 {file_path} 时出错: {str(e)}")
            return None
    
    def read_xlsx(self, file_path):
        """读取 XLSX 文件"""
        if not XLSX_AVAILABLE:
            return None
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"[工作表: {sheet_name}]")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            cell_str = str(cell).strip()
                            if cell_str:
                                row_text.append(cell_str)
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            return "\n".join(text_parts)
        except Exception as e:
            print(f"读取 XLSX 文件 {file_path} 时出错: {str(e)}")
            return None
    
    def read_xls(self, file_path):
        """读取 XLS 文件（旧版 Excel）"""
        if not XLS_AVAILABLE:
            return None
        
        try:
            workbook = xlrd.open_workbook(file_path)
            text_parts = []
            
            for sheet in workbook.sheets():
                text_parts.append(f"[工作表: {sheet.name}]")
                
                for row_idx in range(sheet.nrows):
                    row_text = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        cell_value = cell.value
                        if cell_value:
                            cell_str = str(cell_value).strip()
                            if cell_str:
                                row_text.append(cell_str)
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            return "\n".join(text_parts)
        except Exception as e:
            print(f"读取 XLS 文件 {file_path} 时出错: {str(e)}")
            return None
    
    def read_pptx(self, file_path):
        """读取 PPTX 文件"""
        if not PPTX_AVAILABLE:
            return None
        
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                text_parts.append(f"[幻灯片 {slide_idx}]")
                
                for shape in slide.shapes:
                    # 读取文本内容
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                    
                    # 读取表格
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text_parts.append(" | ".join(row_text))
            
            return "\n".join(text_parts)
        except Exception as e:
            print(f"读取 PPTX 文件 {file_path} 时出错: {str(e)}")
            return None
    
    def read_pcap(self, file_path):
        """读取 PCAP/PCAPNG 流量包文件"""
        if not SCAPY_AVAILABLE:
            return None
        
        try:
            packets = rdpcap(file_path)
            text_parts = []
            packet_count = 0
            
            for packet in packets:
                packet_count += 1
                packet_info = []
                
                # 提取 IP 层信息
                if IP in packet:
                    ip_layer = packet[IP]
                    src_ip = ip_layer.src
                    dst_ip = ip_layer.dst
                    protocol = ip_layer.proto
                    packet_info.append(f"[数据包 #{packet_count}] {src_ip} -> {dst_ip}")
                
                # HTTP 请求
                if HTTP_AVAILABLE and HTTPRequest and HTTPRequest in packet:
                    http_req = packet[HTTPRequest]
                    packet_info.append(f"  [HTTP 请求]")
                    if hasattr(http_req, 'Host'):
                        packet_info.append(f"    Host: {http_req.Host.decode('utf-8', errors='ignore')}")
                    if hasattr(http_req, 'Path'):
                        packet_info.append(f"    Path: {http_req.Path.decode('utf-8', errors='ignore')}")
                    if hasattr(http_req, 'Method'):
                        packet_info.append(f"    Method: {http_req.Method.decode('utf-8', errors='ignore')}")
                    
                    # 提取 HTTP 头部
                    if hasattr(http_req, 'headers'):
                        for key, value in http_req.headers.items():
                            try:
                                key_str = key.decode('utf-8', errors='ignore') if isinstance(key, bytes) else str(key)
                                value_str = value.decode('utf-8', errors='ignore') if isinstance(value, bytes) else str(value)
                                packet_info.append(f"    {key_str}: {value_str}")
                            except:
                                pass
                    
                    # 提取 HTTP 请求体
                    if Raw in packet:
                        try:
                            body = packet[Raw].load
                            if body:
                                body_str = body.decode('utf-8', errors='ignore')
                                if body_str.strip():
                                    packet_info.append(f"    Body: {body_str[:500]}")  # 限制长度
                        except:
                            pass
                
                # HTTP 响应
                elif HTTP_AVAILABLE and HTTPResponse and HTTPResponse in packet:
                    http_resp = packet[HTTPResponse]
                    packet_info.append(f"  [HTTP 响应]")
                    if hasattr(http_resp, 'Status_Code'):
                        packet_info.append(f"    Status Code: {http_resp.Status_Code}")
                    if hasattr(http_resp, 'Reason_Phrase'):
                        try:
                            reason = http_resp.Reason_Phrase
                            if isinstance(reason, bytes):
                                reason = reason.decode('utf-8', errors='ignore')
                            packet_info.append(f"    Reason: {reason}")
                        except:
                            pass
                    
                    # 提取 HTTP 响应头部
                    if hasattr(http_resp, 'headers'):
                        for key, value in http_resp.headers.items():
                            try:
                                key_str = key.decode('utf-8', errors='ignore') if isinstance(key, bytes) else str(key)
                                value_str = value.decode('utf-8', errors='ignore') if isinstance(value, bytes) else str(value)
                                packet_info.append(f"    {key_str}: {value_str}")
                            except:
                                pass
                    
                    # 提取 HTTP 响应体
                    if Raw in packet:
                        try:
                            body = packet[Raw].load
                            if body:
                                body_str = body.decode('utf-8', errors='ignore')
                                if body_str.strip():
                                    packet_info.append(f"    Body: {body_str[:500]}")  # 限制长度
                        except:
                            pass
                
                # DNS 查询
                elif DNS in packet:
                    dns = packet[DNS]
                    packet_info.append(f"  [DNS]")
                    if dns.qr == 0:  # 查询
                        if dns.qd:
                            qname = dns.qd.qname.decode('utf-8', errors='ignore').rstrip('.')
                            qtype = dns.qd.qtype
                            packet_info.append(f"    查询: {qname} (类型: {qtype})")
                    else:  # 响应
                        if dns.an:
                            for i in range(dns.ancount):
                                try:
                                    rr = dns.an[i]
                                    if hasattr(rr, 'rdata'):
                                        rdata = rr.rdata
                                        if isinstance(rdata, bytes):
                                            try:
                                                rdata = rdata.decode('utf-8', errors='ignore')
                                            except:
                                                rdata = str(rdata)
                                        packet_info.append(f"    响应: {rdata}")
                                except:
                                    pass
                
                # TCP/UDP 原始数据（尝试提取文本协议）
                elif Raw in packet and (TCP in packet or UDP in packet):
                    try:
                        raw_data = packet[Raw].load
                        if len(raw_data) > 0:
                            # 尝试多种编码
                            text_data = None
                            for encoding in ['utf-8', 'gbk', 'latin-1']:
                                try:
                                    text_data = raw_data.decode(encoding, errors='ignore')
                                    # 检查是否包含可打印字符
                                    if text_data and any(c.isprintable() or c in '\r\n\t' for c in text_data[:200]):
                                        break
                                except:
                                    continue
                            
                            if text_data:
                                # 检查是否像 HTTP、FTP、SMTP 等文本协议
                                text_upper = text_data[:500].upper()
                                protocol_keywords = ['HTTP', 'GET ', 'POST ', 'PUT ', 'DELETE ', 'HEAD ', 'OPTIONS ',
                                                   'FTP', 'SMTP', 'POP3', 'IMAP', '220 ', '250 ', '331 ', '334 ']
                                
                                if any(keyword in text_upper for keyword in protocol_keywords):
                                    lines = text_data.split('\n')[:30]  # 取前30行
                                    packet_info.append(f"  [协议数据]")
                                    for line in lines:
                                        line_stripped = line.strip()
                                        if line_stripped and len(line_stripped) > 0:
                                            # 限制每行长度，避免过长
                                            display_line = line_stripped[:300]
                                            packet_info.append(f"    {display_line}")
                    except Exception as e:
                        # 静默失败，不输出错误
                        pass
                
                # 如果有信息，添加到结果
                if packet_info:
                    text_parts.extend(packet_info)
                    text_parts.append("")  # 空行分隔
            
            if text_parts:
                header = f"[流量包分析] 共 {packet_count} 个数据包\n" + "=" * 50 + "\n"
                return header + "\n".join(text_parts)
            else:
                return f"[流量包分析] 共 {packet_count} 个数据包，但未提取到文本内容"
        except Exception as e:
            print(f"读取流量包文件 {file_path} 时出错: {str(e)}")
            return None
    
    def read_evtx(self, file_path):
        """读取 Windows 事件日志文件（.evtx）"""
        if not EVTX_AVAILABLE:
            return None
        
        try:
            text_parts = []
            event_count = 0
            
            with Evtx(file_path) as log:
                for record in log.records():
                    event_count += 1
                    try:
                        # 解析 XML 事件数据
                        xml_data = record.xml()
                        root = ET.fromstring(xml_data)
                        
                        # 提取事件信息
                        event_info = []
                        event_info.append(f"[事件 #{event_count}]")
                        
                        # 提取系统信息
                        system = root.find('.//System')
                        if system is not None:
                            # 事件 ID
                            event_id_elem = system.find('EventID')
                            if event_id_elem is not None:
                                event_id = event_id_elem.text if event_id_elem.text else event_id_elem.get('Value', '')
                                event_info.append(f"  事件ID: {event_id}")
                            
                            # 事件级别
                            level_elem = system.find('Level')
                            if level_elem is not None:
                                level = level_elem.text if level_elem.text else level_elem.get('Value', '')
                                level_map = {'1': '关键', '2': '错误', '3': '警告', '4': '信息', '5': '详细'}
                                level_text = level_map.get(level, level)
                                event_info.append(f"  级别: {level_text}")
                            
                            # 时间戳
                            time_created = system.find('TimeCreated')
                            if time_created is not None:
                                timestamp = time_created.get('SystemTime', '')
                                if timestamp:
                                    event_info.append(f"  时间: {timestamp}")
                            
                            # 计算机名
                            computer = system.find('Computer')
                            if computer is not None and computer.text:
                                event_info.append(f"  计算机: {computer.text}")
                            
                            # 提供程序
                            provider = system.find('Provider')
                            if provider is not None:
                                provider_name = provider.get('Name', '')
                                if provider_name:
                                    event_info.append(f"  提供程序: {provider_name}")
                            
                            # 任务类别
                            task = system.find('Task')
                            if task is not None:
                                task_text = task.text if task.text else task.get('Value', '')
                                if task_text:
                                    event_info.append(f"  任务: {task_text}")
                            
                            # 关键字
                            keywords = system.find('Keywords')
                            if keywords is not None:
                                keyword_text = keywords.text if keywords.text else keywords.get('Value', '')
                                if keyword_text:
                                    event_info.append(f"  关键字: {keyword_text}")
                        
                        # 提取事件数据
                        event_data = root.find('.//EventData')
                        if event_data is not None:
                            data_items = event_data.findall('Data')
                            if data_items:
                                event_info.append("  事件数据:")
                                for data_item in data_items:
                                    name = data_item.get('Name', '')
                                    value = data_item.text if data_item.text else ''
                                    if value:
                                        if name:
                                            event_info.append(f"    {name}: {value[:500]}")
                                        else:
                                            event_info.append(f"    {value[:500]}")
                        
                        # 提取用户数据
                        user_data = root.find('.//UserData')
                        if user_data is not None:
                            # 提取所有文本内容
                            user_text = ET.tostring(user_data, encoding='unicode', method='text')
                            if user_text and user_text.strip():
                                event_info.append(f"  用户数据: {user_text.strip()[:500]}")
                        
                        # 提取消息/描述
                        message = root.find('.//Message')
                        if message is not None and message.text:
                            event_info.append(f"  消息: {message.text[:500]}")
                        
                        # 提取所有文本节点（作为备用）
                        all_text = ET.tostring(root, encoding='unicode', method='text')
                        if all_text and all_text.strip():
                            # 只提取前1000个字符，避免过长
                            text_preview = ' '.join(all_text.split()[:50])  # 取前50个词
                            if len(text_preview) > 0:
                                event_info.append(f"  内容预览: {text_preview[:500]}")
                        
                        # 如果有信息，添加到结果
                        if len(event_info) > 1:  # 至少有事件编号
                            text_parts.extend(event_info)
                            text_parts.append("")  # 空行分隔
                    
                    except Exception as e:
                        # 单个事件解析失败，继续处理下一个
                        print(f"解析事件 #{event_count} 时出错: {str(e)}")
                        continue
            
            if text_parts:
                header = f"[Windows 事件日志分析] 共 {event_count} 个事件\n" + "=" * 50 + "\n"
                return header + "\n".join(text_parts)
            else:
                return f"[Windows 事件日志分析] 共 {event_count} 个事件，但未提取到文本内容"
        except Exception as e:
            print(f"读取 Windows 事件日志文件 {file_path} 时出错: {str(e)}")
            return None
    
    def read_file(self, file_path):
        """根据文件类型自动选择读取方法"""
        ext = self.get_file_extension(file_path)
        
        # 纯文本文件
        if self.is_text_file(file_path):
            return self.read_text_file(file_path)
        
        # DOCX 文件
        elif ext == '.docx':
            result = self.read_docx(file_path)
            if result is None and DOCX_AVAILABLE is False:
                print(f"提示: 无法读取 {file_path}，请安装 python-docx 库: pip install python-docx")
            return result
        
        # XLSX 文件
        elif ext == '.xlsx':
            result = self.read_xlsx(file_path)
            if result is None and XLSX_AVAILABLE is False:
                print(f"提示: 无法读取 {file_path}，请安装 openpyxl 库: pip install openpyxl")
            return result
        
        # XLS 文件
        elif ext == '.xls':
            result = self.read_xls(file_path)
            if result is None and XLS_AVAILABLE is False:
                print(f"提示: 无法读取 {file_path}，请安装 xlrd 库: pip install xlrd")
            return result
        
        # PPTX 文件
        elif ext == '.pptx':
            result = self.read_pptx(file_path)
            if result is None and PPTX_AVAILABLE is False:
                print(f"提示: 无法读取 {file_path}，请安装 python-pptx 库: pip install python-pptx")
            return result
        
        # 旧版 Office 格式（.doc, .ppt）- 需要额外处理，暂时尝试作为文本读取
        elif ext in ['.doc', '.ppt']:
            print(f"提示: {ext} 格式（旧版 Office）暂不支持直接读取，将尝试作为文本文件处理")
            return self.read_text_file(file_path)
        
        # PDF 文件 - 需要额外库支持
        elif ext == '.pdf':
            print(f"提示: PDF 文件暂不支持直接读取，将尝试作为文本文件处理")
            print(f"如需支持 PDF，可安装 PyPDF2 或 pdfplumber 库")
            return self.read_text_file(file_path)
        
        # 流量包文件
        elif ext in ['.pcap', '.pcapng', '.cap']:
            result = self.read_pcap(file_path)
            if result is None and SCAPY_AVAILABLE is False:
                print(f"提示: 无法读取 {file_path}，请安装 scapy 库: pip install scapy")
            return result
        
        # Windows 事件日志文件
        elif ext == '.evtx':
            result = self.read_evtx(file_path)
            if result is None and EVTX_AVAILABLE is False:
                print(f"提示: 无法读取 {file_path}，请安装 python-evtx 库: pip install python-evtx")
            return result
        
        # 其他格式，尝试作为文本文件读取
        else:
            return self.read_text_file(file_path)

